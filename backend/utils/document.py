"""
Document ingestion utilities with format-aware parsing and semantic structure retention.
"""

from __future__ import annotations

from collections import Counter
from dataclasses import dataclass, field
import io
import re

from pypdf import PdfReader


TEXT_EXTENSIONS = {"txt", "md", "csv", "json", "py", "js", "ts", "tsx", "jsx"}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "webp", "bmp", "tiff", "tif"}
SUPPORTED_EXTENSIONS = {
    "pdf",
    "docx",
    "pptx",
    "xlsx",
    *TEXT_EXTENSIONS,
    *IMAGE_EXTENSIONS,
}


@dataclass
class DocumentSegment:
    text: str
    kind: str
    section_title: str | None = None
    page_number: int | None = None
    slide_number: int | None = None
    sheet_name: str | None = None
    level: int | None = None


@dataclass
class ParsedDocument:
    filename: str
    file_type: str
    title: str
    segments: list[DocumentSegment] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        return "\n\n".join(segment.text for segment in self.segments if segment.text.strip())


@dataclass
class DocumentProfile:
    category: str
    inferred_intent: str
    primary_subject: str
    top_headings: list[str] = field(default_factory=list)
    focus_terms: list[str] = field(default_factory=list)
    named_entities: list[str] = field(default_factory=list)
    prompt_context: list[str] = field(default_factory=list)

    def as_dict(self) -> dict:
        return {
            "category": self.category,
            "inferred_intent": self.inferred_intent,
            "primary_subject": self.primary_subject,
            "top_headings": self.top_headings,
            "focus_terms": self.focus_terms,
            "named_entities": self.named_entities,
            "prompt_context": self.prompt_context,
        }


STOPWORDS = {
    "a",
    "an",
    "and",
    "are",
    "as",
    "at",
    "be",
    "by",
    "for",
    "from",
    "in",
    "into",
    "is",
    "it",
    "of",
    "on",
    "or",
    "that",
    "the",
    "this",
    "to",
    "with",
    "using",
    "use",
    "via",
    "your",
    "their",
    "these",
    "those",
    "data",
    "document",
    "page",
    "slide",
    "section",
}


def _dedupe_preserve(values: list[str]) -> list[str]:
    seen: set[str] = set()
    ordered: list[str] = []
    for value in values:
        key = value.strip().lower()
        if not key or key in seen:
            continue
        seen.add(key)
        ordered.append(value.strip())
    return ordered


def _extract_heading_candidates(parsed: ParsedDocument) -> list[str]:
    headings = [
        segment.section_title or segment.text
        for segment in parsed.segments
        if segment.kind in {"heading", "slide_title"} or (segment.level is not None and segment.level <= 1)
    ]
    if not headings:
        headings = [segment.section_title or "" for segment in parsed.segments[:8]]
    cleaned = [normalize_text(item) for item in headings if item]
    return _dedupe_preserve(cleaned)[:4]


def _extract_focus_terms(parsed: ParsedDocument) -> list[str]:
    text = normalize_text(parsed.full_text[:12000]).lower()
    terms = re.findall(r"\b[a-z][a-z0-9\-\+]{2,}\b", text)
    counts = Counter(term for term in terms if term not in STOPWORDS and not term.isdigit())
    ranked = [term for term, _ in counts.most_common(8)]
    return ranked[:5]


def _extract_named_entities(parsed: ParsedDocument) -> list[str]:
    candidates: list[str] = [parsed.title]
    candidates.extend(_extract_heading_candidates(parsed))

    for segment in parsed.segments[:10]:
        sample = normalize_text(segment.text).replace("\n", " ")
        if not sample:
            continue
        sentence = sample.split(". ", 1)[0][:140]
        candidates.extend(
            re.findall(
                r"\b(?:[A-Z][a-z]+(?:\s+[A-Z][a-zA-Z0-9&\-]+){0,2}|[A-Z]{2,}(?:\s+[A-Z]{2,}){0,1})\b",
                sentence,
            )
        )

    filtered = []
    for item in candidates:
        value = normalize_text(item).replace("\n", " ")
        if not (2 < len(value) < 48):
            continue
        if value.lower() in {"page", "section", "slide", "table", "experience", "projects", "education"}:
            continue
        filtered.append(value)
    return _dedupe_preserve(filtered)[:4]


def _build_prompt_context(parsed: ParsedDocument, headings: list[str], focus_terms: list[str]) -> list[str]:
    context: list[str] = []
    if headings:
        context.append(headings[0])
    first_substantive = next(
        (
            normalize_text(segment.text)
            for segment in parsed.segments
            if len(normalize_text(segment.text)) > 40 and segment.kind not in {"heading", "slide_title"}
        ),
        "",
    )
    if first_substantive:
        context.append(first_substantive[:120])
    if focus_terms:
        context.append(", ".join(focus_terms[:3]))
    return _dedupe_preserve(context)[:3]


def infer_document_category(parsed: ParsedDocument) -> str:
    haystack = f"{parsed.filename} {parsed.title} {parsed.full_text[:5000]}".lower()
    if any(token in haystack for token in ["resume", "curriculum vitae", "experience", "education", "skills"]):
        return "resume"
    if any(token in haystack for token in ["abstract", "methodology", "results", "references", "experiment"]):
        return "research"
    if any(token in haystack for token in ["roadmap", "milestone", "timeline", "deliverable", "project"]):
        return "project"
    if any(token in haystack for token in ["budget", "forecast", "revenue", "financial", "report"]):
        return "report"
    return "general"


def build_document_profile(parsed: ParsedDocument) -> DocumentProfile:
    category = infer_document_category(parsed)
    headings = _extract_heading_candidates(parsed)
    focus_terms = _extract_focus_terms(parsed)
    entities = _extract_named_entities(parsed)
    prompt_context = _build_prompt_context(parsed, headings, focus_terms)

    primary_subject = headings[0] if headings else (entities[0] if entities else parsed.title)
    if category == "resume":
        inferred_intent = "candidate evaluation"
    elif category == "research":
        inferred_intent = "paper analysis"
    elif category == "project":
        inferred_intent = "project planning review"
    elif category == "report":
        inferred_intent = "report synthesis"
    else:
        inferred_intent = "document analysis"

    return DocumentProfile(
        category=category,
        inferred_intent=inferred_intent,
        primary_subject=primary_subject,
        top_headings=headings,
        focus_terms=focus_terms,
        named_entities=entities,
        prompt_context=prompt_context,
    )


def build_suggested_prompts(parsed: ParsedDocument, profile: DocumentProfile | None = None) -> list[str]:
    profile = profile or build_document_profile(parsed)
    category = profile.category
    subject = profile.primary_subject
    heading = profile.top_headings[1] if len(profile.top_headings) > 1 else (profile.top_headings[0] if profile.top_headings else subject)
    term_a = profile.focus_terms[0] if profile.focus_terms else "key themes"
    term_b = profile.focus_terms[1] if len(profile.focus_terms) > 1 else "technical details"
    parsed_title = parsed.title.lower()
    related_entities = [
        item
        for item in profile.named_entities
        if item.lower() not in {subject.lower(), parsed_title}
    ]
    entity_a = subject
    entity_b = related_entities[0] if related_entities else heading

    if category == "resume":
        return [
            f"Summarize this candidate for a backend, AI, or systems role using evidence from {entity_a}.",
            f"What engineering signals, shipped projects, and measurable outcomes stand out across {entity_a} and {entity_b}?",
            f"Which experience best supports backend, AI, or systems credibility, especially around {term_a} and {term_b}?",
        ]
    if category == "research":
        return [
            f"What is the core contribution or thesis of {entity_a}?",
            f"Summarize the methodology, experiments, and results tied to {term_a} and {term_b}.",
            f"What limitations, assumptions, or future work are described in {heading}?",
        ]
    if category == "project":
        return [
            f"What milestones, deliverables, or execution phases are defined for {entity_a}?",
            f"What risks, blockers, or dependencies are called out around {term_a} and {term_b}?",
            f"Extract action items, owners, or next steps mentioned in {heading} or related to {entity_b}.",
        ]
    if category == "report":
        return [
            f"Summarize the key findings or takeaways from {entity_a}.",
            f"What trends, risks, or recommendations are associated with {term_a} and {term_b}?",
            f"Which metrics, comparisons, or conclusions in {heading} matter most?",
        ]
    prompts = [
        f"Summarize the most important points in {entity_a}.",
        f"What facts, decisions, or findings are emphasized around {term_a}?",
        f"Which section or evidence in {heading} is most relevant to this document's purpose?",
    ]
    if profile.named_entities:
        prompts.append(f"What roles, organizations, or named entities like {profile.named_entities[0]} are important here?")
    return _dedupe_preserve(prompts)[:4]


def get_file_extension(filename: str) -> str:
    return filename.lower().rsplit(".", 1)[-1] if "." in filename else ""


def is_image_file(filename: str) -> bool:
    return get_file_extension(filename) in IMAGE_EXTENSIONS


def normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _title_from_filename(filename: str) -> str:
    base = filename.rsplit(".", 1)[0]
    return base.replace("_", " ").replace("-", " ").strip() or filename


def _add_segment(
    segments: list[DocumentSegment],
    *,
    text: str,
    kind: str,
    section_title: str | None = None,
    page_number: int | None = None,
    slide_number: int | None = None,
    sheet_name: str | None = None,
    level: int | None = None,
) -> None:
    normalized = normalize_text(text)
    if not normalized:
        return
    segments.append(
        DocumentSegment(
            text=normalized,
            kind=kind,
            section_title=section_title,
            page_number=page_number,
            slide_number=slide_number,
            sheet_name=sheet_name,
            level=level,
        )
    )


def extract_pdf_segments(file_bytes: bytes, filename: str) -> ParsedDocument:
    title = _title_from_filename(filename)
    segments: list[DocumentSegment] = []

    try:
        import fitz  # PyMuPDF

        document = fitz.open(stream=file_bytes, filetype="pdf")
        for page_index, page in enumerate(document, start=1):
            blocks = page.get_text("blocks")
            page_segments = 0
            for block in sorted(blocks, key=lambda item: (item[1], item[0])):
                text = block[4]
                if normalize_text(text):
                    _add_segment(
                        segments,
                        text=text,
                        kind="paragraph",
                        section_title=f"Page {page_index}",
                        page_number=page_index,
                    )
                    page_segments += 1
            if page_segments == 0:
                fallback_text = normalize_text(page.get_text("text"))
                if fallback_text:
                    _add_segment(
                        segments,
                        text=fallback_text,
                        kind="paragraph",
                        section_title=f"Page {page_index}",
                        page_number=page_index,
                    )
        document.close()
    except ImportError:
        reader = PdfReader(io.BytesIO(file_bytes))
        for page_index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            _add_segment(
                segments,
                text=text,
                kind="paragraph",
                section_title=f"Page {page_index}",
                page_number=page_index,
            )

    return ParsedDocument(filename=filename, file_type="pdf", title=title, segments=segments)


def extract_docx_segments(file_bytes: bytes, filename: str) -> ParsedDocument:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("DOCX support requires 'python-docx' to be installed.") from exc

    title = _title_from_filename(filename)
    doc = Document(io.BytesIO(file_bytes))
    segments: list[DocumentSegment] = []
    current_heading = title

    for paragraph in doc.paragraphs:
        text = normalize_text(paragraph.text)
        if not text:
            continue

        style_name = (paragraph.style.name if paragraph.style else "").lower()
        if "heading" in style_name or style_name in {"title", "subtitle"}:
            current_heading = text
            level_match = re.search(r"(\d+)", style_name)
            level = int(level_match.group(1)) if level_match else 1
            _add_segment(
                segments,
                text=text,
                kind="heading",
                section_title=current_heading,
                level=level,
            )
        else:
            _add_segment(
                segments,
                text=text,
                kind="paragraph",
                section_title=current_heading,
            )

    for table_index, table in enumerate(doc.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            cells = [normalize_text(cell.text) for cell in row.cells]
            values = [cell for cell in cells if cell]
            if values:
                rows.append(" | ".join(values))
        if rows:
            _add_segment(
                segments,
                text="\n".join(rows),
                kind="table",
                section_title=f"{current_heading} / Table {table_index}",
            )

    return ParsedDocument(filename=filename, file_type="docx", title=title, segments=segments)


def extract_pptx_segments(file_bytes: bytes, filename: str) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PPTX support requires 'python-pptx' to be installed.") from exc

    title = _title_from_filename(filename)
    deck = Presentation(io.BytesIO(file_bytes))
    segments: list[DocumentSegment] = []

    for slide_index, slide in enumerate(deck.slides, start=1):
        slide_title = title
        if getattr(slide.shapes, "title", None) and slide.shapes.title.text:
            slide_title = normalize_text(slide.shapes.title.text) or slide_title

        _add_segment(
            segments,
            text=slide_title,
            kind="slide_title",
            section_title=slide_title,
            slide_number=slide_index,
            level=1,
        )

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                text = getattr(shape, "text", None)
                if text and normalize_text(text) and normalize_text(text) != slide_title:
                    _add_segment(
                        segments,
                        text=text,
                        kind="paragraph",
                        section_title=slide_title,
                        slide_number=slide_index,
                    )
                continue

            for paragraph in shape.text_frame.paragraphs:
                runs = [run.text for run in paragraph.runs if run.text]
                text = normalize_text("".join(runs) if runs else paragraph.text)
                if not text or text == slide_title:
                    continue
                kind = "bullet" if paragraph.level > 0 else "paragraph"
                _add_segment(
                    segments,
                    text=text,
                    kind=kind,
                    section_title=slide_title,
                    slide_number=slide_index,
                    level=paragraph.level,
                )

    return ParsedDocument(filename=filename, file_type="pptx", title=title, segments=segments)


def extract_xlsx_segments(file_bytes: bytes, filename: str) -> ParsedDocument:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("XLSX support requires 'openpyxl' to be installed.") from exc

    title = _title_from_filename(filename)
    workbook = load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
    segments: list[DocumentSegment] = []

    for worksheet in workbook.worksheets:
        _add_segment(
            segments,
            text=worksheet.title,
            kind="heading",
            section_title=worksheet.title,
            sheet_name=worksheet.title,
            level=1,
        )
        for row in worksheet.iter_rows(values_only=True):
            values = [normalize_text(str(value)) for value in row if value is not None]
            filtered = [value for value in values if value]
            if filtered:
                _add_segment(
                    segments,
                    text=" | ".join(filtered),
                    kind="table_row",
                    section_title=worksheet.title,
                    sheet_name=worksheet.title,
                )

    return ParsedDocument(filename=filename, file_type="xlsx", title=title, segments=segments)


def extract_text_segments(file_bytes: bytes, filename: str) -> ParsedDocument:
    title = _title_from_filename(filename)
    decoded = file_bytes.decode("utf-8", errors="replace")
    blocks = [normalize_text(block) for block in decoded.split("\n\n")]
    segments: list[DocumentSegment] = []
    current_heading = title

    for block in blocks:
        if not block:
            continue
        lines = [normalize_text(line) for line in block.split("\n") if normalize_text(line)]
        if not lines:
            continue

        first_line = lines[0]
        remaining = "\n".join(lines[1:]).strip()
        looks_like_heading = (
            len(first_line) <= 64
            and len(first_line.split()) <= 8
            and not first_line.endswith(".")
            and first_line.lower() != title.lower()
        )

        if looks_like_heading:
            current_heading = first_line
            _add_segment(segments, text=first_line, kind="heading", section_title=current_heading, level=1)
            if remaining:
                _add_segment(segments, text=remaining, kind="paragraph", section_title=current_heading)
        else:
            _add_segment(segments, text=block, kind="paragraph", section_title=current_heading)

    return ParsedDocument(filename=filename, file_type=get_file_extension(filename), title=title, segments=segments)


def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    return extract_structured_document(file_bytes, filename).full_text


def extract_structured_document(file_bytes: bytes, filename: str) -> ParsedDocument:
    ext = get_file_extension(filename)

    if ext == "pdf":
        return extract_pdf_segments(file_bytes, filename)
    if ext == "docx":
        return extract_docx_segments(file_bytes, filename)
    if ext == "pptx":
        return extract_pptx_segments(file_bytes, filename)
    if ext == "xlsx":
        return extract_xlsx_segments(file_bytes, filename)
    if ext in IMAGE_EXTENSIONS:
        raise RuntimeError("Image files must be routed through the vision OCR pipeline.")
    if ext in TEXT_EXTENSIONS:
        return extract_text_segments(file_bytes, filename)

    raise RuntimeError(
        f"Unsupported file type '.{ext}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
    )
